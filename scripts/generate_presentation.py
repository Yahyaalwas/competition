#!/usr/bin/env python3
"""
GCC AI Intelligence Platform — Competition Presentation Generator
Produces a clean, professional 10-slide deck matching competition requirements.

Usage:
    python scripts/generate_presentation.py
Output:
    presentation/GCC_AI_Intelligence_Platform.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x4F, 0x72)
NAVY2  = RGBColor(0x0D, 0x2A, 0x42)
GOLD   = RGBColor(0xC3, 0x9B, 0x4E)
GOLD2  = RGBColor(0xE8, 0xC9, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x09, 0x13, 0x20)
GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
RED    = RGBColor(0xA9, 0x32, 0x26)
AMBER  = RGBColor(0xC0, 0x78, 0x20)
TEAL   = RGBColor(0x14, 0x8F, 0x77)
BLUE   = RGBColor(0x28, 0x74, 0xA6)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MGRAY  = RGBColor(0x88, 0x88, 0x88)

SW = Inches(13.33)
SH = Inches(7.5)
MSO_RECT = 1


# ── Core drawing helpers ──────────────────────────────────────────────────────
def _prs():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    sp = slide.shapes.add_shape(MSO_RECT, l, t, w, h)
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = line
        sp.line.width = lw
    else:
        sp.line.fill.background()
    return sp

def _txt(slide, text, l, t, w, h,
         size=Pt(11), bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text         = text
    r.font.size    = size
    r.font.bold    = bold
    r.font.color.rgb = color
    r.font.italic  = italic
    return tb

def _para(slide, lines, l, t, w, h, default_size=Pt(11),
          default_color=WHITE, align=PP_ALIGN.LEFT, spacing=Pt(4)):
    """Multi-paragraph textbox. lines = list of str or (text,bold,size,color)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, False, default_size, default_color)
        text, bold, sz, col = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = spacing
        r = p.add_run()
        r.text = text
        r.font.size  = sz
        r.font.bold  = bold
        r.font.color.rgb = col
    return tb


# ── Branded layout helpers ────────────────────────────────────────────────────
def _header(slide, title_en, title_ar="", tag=""):
    """Full-width dark header bar with gold underline."""
    _rect(slide, Inches(0), Inches(0), SW, Inches(1.22), fill=DARK)
    _rect(slide, Inches(0), Inches(1.22), SW, Inches(0.06), fill=GOLD)
    _txt(slide, title_en,
         Inches(0.5), Inches(0.1), Inches(10.5), Inches(0.72),
         size=Pt(28), bold=True, color=WHITE)
    if title_ar:
        _txt(slide, title_ar,
             Inches(0.5), Inches(0.76), Inches(10), Inches(0.38),
             size=Pt(13), color=GOLD2, italic=True)
    if tag:
        _rect(slide, Inches(11.2), Inches(0.3), Inches(1.95), Inches(0.55), fill=GOLD)
        _txt(slide, tag, Inches(11.2), Inches(0.3), Inches(1.95), Inches(0.55),
             size=Pt(9), bold=True, color=DARK, align=PP_ALIGN.CENTER)

def _footer(slide):
    _rect(slide, Inches(0), Inches(7.2), SW, Inches(0.3), fill=DARK)
    _txt(slide,
         "GCC AI Intelligence Platform  ·  Yahya Alwashahi  ·  World Bank Open Data  ·  ycaward@msy.gov.qa",
         Inches(0.5), Inches(7.22), Inches(12.5), Inches(0.25),
         size=Pt(7.5), color=GOLD, align=PP_ALIGN.CENTER)

def _section_tag(slide, en, ar, x, y, color=NAVY):
    """Small section label — EN above, AR below."""
    _rect(slide, x, y, Inches(2.4), Inches(0.52), fill=color)
    _txt(slide, en, x + Inches(0.1), y + Inches(0.04), Inches(2.2), Inches(0.24),
         size=Pt(8), bold=True, color=WHITE)
    _txt(slide, ar, x + Inches(0.1), y + Inches(0.28), Inches(2.2), Inches(0.22),
         size=Pt(8), color=GOLD2)

def _stat_block(slide, val, lbl_en, lbl_ar, x, y, w=Inches(2.55), h=Inches(1.55),
                bg=NAVY, vc=GOLD, lc=WHITE):
    _rect(slide, x, y, w, h, fill=bg)
    _txt(slide, val, x, y + Inches(0.12), w, Inches(0.72),
         size=Pt(38), bold=True, color=vc, align=PP_ALIGN.CENTER)
    _txt(slide, lbl_en, x, y + Inches(0.82), w, Inches(0.3),
         size=Pt(9), bold=True, color=lc, align=PP_ALIGN.CENTER)
    _txt(slide, lbl_ar, x, y + Inches(1.12), w, Inches(0.28),
         size=Pt(8.5), color=GOLD2, align=PP_ALIGN.CENTER)

def _bullet_list(slide, bullets, l, t, w, h,
                 size=Pt(11), color=DARK, marker="▸  "):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if isinstance(b, tuple):
            txt, sz, col = b
        else:
            txt, sz, col = b, size, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = marker + txt
        r.font.size  = sz
        r.font.color.rgb = col

def _screenshot_box(slide, l, t, w, h, label_en, label_ar, tip=""):
    _rect(slide, l, t, w, h, fill=LGRAY, line=NAVY, lw=Pt(1.5))
    _rect(slide, l, t, w, Inches(0.4), fill=NAVY)
    _txt(slide, "📸  " + label_en, l + Inches(0.1), t + Inches(0.06), w - Inches(0.2), Inches(0.3),
         size=Pt(8.5), bold=True, color=WHITE)
    _txt(slide, label_ar,
         l + Inches(0.1), t + h/2 - Inches(0.15), w - Inches(0.2), Inches(0.3),
         size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER, italic=True)
    if tip:
        _txt(slide, tip, l + Inches(0.1), t + h - Inches(0.38), w - Inches(0.2), Inches(0.3),
             size=Pt(7.5), color=GOLD, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ══════════════════════════════════════════════════════════════════════════════
def slide_01_title(prs):
    s = _blank(prs)
    _bg(s, NAVY2)
    # Gold left accent bar
    _rect(s, Inches(0), Inches(0), Inches(0.22), SH, fill=GOLD)
    # Top dark band
    _rect(s, Inches(0.22), Inches(0), SW - Inches(0.22), Inches(0.12), fill=DARK)

    # Globe icon area
    _rect(s, Inches(0.5), Inches(0.5), Inches(1.0), Inches(1.0), fill=GOLD)
    _txt(s, "🌍", Inches(0.5), Inches(0.52), Inches(1.0), Inches(1.0),
         size=Pt(32), align=PP_ALIGN.CENTER)

    # English title
    _txt(s, "GCC AI Intelligence Platform",
         Inches(0.5), Inches(1.75), Inches(9.5), Inches(1.2),
         size=Pt(48), bold=True, color=WHITE)

    # Gold rule
    _rect(s, Inches(0.5), Inches(3.08), Inches(7.0), Inches(0.07), fill=GOLD)

    # Arabic title
    _txt(s, "منصة الذكاء الاصطناعي لتوظيف الشباب الخليجي",
         Inches(0.5), Inches(3.22), Inches(10), Inches(0.7),
         size=Pt(20), color=GOLD2, italic=True)

    # Tagline EN
    _txt(s, "AI-Powered Youth Employment Forecasting & Policy Intelligence for the Gulf Cooperation Council",
         Inches(0.5), Inches(4.05), Inches(9.5), Inches(0.65),
         size=Pt(14), color=WHITE, italic=False)

    # Tagline AR
    _txt(s, "تنبؤ بتوظيف الشباب ودعم قرارات السياسات الخليجية بالذكاء الاصطناعي",
         Inches(0.5), Inches(4.72), Inches(9.5), Inches(0.5),
         size=Pt(12), color=GOLD2, italic=True)

    # Bottom stats strip
    _rect(s, Inches(0.22), Inches(5.55), SW - Inches(0.22), Inches(1.95), fill=DARK)
    stats = [
        ("6", "GCC Nations", "دول خليجية"),
        ("5", "Indicators", "مؤشرات"),
        ("15 Yrs", "Real Data", "بيانات حقيقية"),
        ("6", "AI Models", "نماذج ذكاء"),
        ("8", "Scenarios", "سيناريوهات"),
        ("EN+AR", "Bilingual", "ثنائي اللغة"),
    ]
    for i, (val, en, ar) in enumerate(stats):
        x = Inches(0.5 + i * 2.1)
        _txt(s, val,   x, Inches(5.68), Inches(2.0), Inches(0.72),
             size=Pt(30), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _txt(s, en,    x, Inches(6.36), Inches(2.0), Inches(0.28),
             size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, ar,    x, Inches(6.64), Inches(2.0), Inches(0.26),
             size=Pt(8), color=GOLD2, align=PP_ALIGN.CENTER)

    # Participant + live URL box
    _rect(s, Inches(10.2), Inches(1.75), Inches(2.9), Inches(3.5), fill=NAVY)
    _txt(s, "SUBMITTED BY", Inches(10.3), Inches(1.85), Inches(2.7), Inches(0.3),
         size=Pt(7.5), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _txt(s, "Yahya Alwashahi", Inches(10.3), Inches(2.18), Inches(2.7), Inches(0.5),
         size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(s, Inches(10.4), Inches(2.75), Inches(2.5), Inches(0.04), fill=GOLD)
    _txt(s, "🔗 LIVE PLATFORM", Inches(10.3), Inches(2.9), Inches(2.7), Inches(0.28),
         size=Pt(8), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _txt(s, "competition-efvjqtcbnxcvt3te6fqkaw\n.streamlit.app",
         Inches(10.3), Inches(3.22), Inches(2.7), Inches(0.7),
         size=Pt(7.5), color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    _txt(s, "Data: World Bank Open Data API v2",
         Inches(10.3), Inches(4.05), Inches(2.7), Inches(0.28),
         size=Pt(7), color=MGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
def slide_02_problem(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "The Challenge: GCC Youth Unemployment",
            "التحدي: بطالة الشباب في دول مجلس التعاون",
            tag="المشكلة\nPROBLEM")
    _footer(s)

    # Left panel — the numbers
    _rect(s, Inches(0.35), Inches(1.42), Inches(4.8), Inches(5.5), fill=DARK)
    _txt(s, "📊  THE CHALLENGE IN NUMBERS",
         Inches(0.55), Inches(1.55), Inches(4.4), Inches(0.35),
         size=Pt(9.5), bold=True, color=GOLD)
    _txt(s, "التحدي بالأرقام",
         Inches(0.55), Inches(1.88), Inches(4.4), Inches(0.28),
         size=Pt(9), color=GOLD2, italic=True)

    stats = [
        ("30%+",  "Youth unemployment in some GCC states",     "معدل بطالة الشباب في بعض دول المجلس"),
        ("3.5%",  "Lowest rate — UAE & Qatar",                  "أدنى معدل — الإمارات وقطر"),
        ("60M",   "People in the GCC workforce ecosystem",      "شخص في منظومة العمل الخليجية"),
        ("#1",    "Strategic priority for all 6 national visions","الأولوية الاستراتيجية لجميع رؤى المجلس الست"),
    ]
    for i, (val, en, ar) in enumerate(stats):
        y = Inches(2.3 + i * 1.1)
        _rect(s, Inches(0.45), y, Inches(4.6), Inches(0.9), fill=NAVY)
        _txt(s, val,  Inches(0.55), y + Inches(0.06), Inches(1.2), Inches(0.55),
             size=Pt(26), bold=True, color=GOLD)
        _txt(s, en,   Inches(1.85), y + Inches(0.08), Inches(3.1), Inches(0.32),
             size=Pt(9.5), bold=True, color=WHITE)
        _txt(s, ar,   Inches(1.85), y + Inches(0.42), Inches(3.1), Inches(0.35),
             size=Pt(8.5), color=GOLD2)

    # Right panel — the 5 gaps
    _rect(s, Inches(5.45), Inches(1.42), Inches(7.55), Inches(5.5), fill=WHITE, line=NAVY, lw=Pt(1.5))
    _rect(s, Inches(5.45), Inches(1.42), Inches(7.55), Inches(0.5), fill=NAVY)
    _txt(s, "⚠  THE POLICY INTELLIGENCE GAP  —  الفجوة في الذكاء السياسي",
         Inches(5.65), Inches(1.48), Inches(7.2), Inches(0.37),
         size=Pt(10), bold=True, color=WHITE)

    gaps = [
        ("1", RED,   "Fragmented Data",          "بيانات مشتتة عبر جهات متعددة"),
        ("2", RED,   "No Regional Benchmark",    "لا توجد مقارنة إقليمية موحدة"),
        ("3", RED,   "No Forward Intelligence",  "لا توجد رؤية تنبؤية مستقبلية"),
        ("4", RED,   "No Scenario Modelling",    "لا توجد محاكاة لتأثير السياسات"),
        ("5", RED,   "No Bilingual Reports",     "لا توجد تقارير ثنائية اللغة للوزارات"),
    ]
    for i, (num, col, en, ar) in enumerate(gaps):
        y = Inches(2.08 + i * 0.95)
        _rect(s, Inches(5.55), y, Inches(0.5), Inches(0.72), fill=col)
        _txt(s, num,  Inches(5.55), y, Inches(0.5), Inches(0.72),
             size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, en,   Inches(6.15), y + Inches(0.04), Inches(6.6), Inches(0.33),
             size=Pt(11), bold=True, color=DARK)
        _txt(s, ar,   Inches(6.15), y + Inches(0.38), Inches(6.6), Inches(0.3),
             size=Pt(9.5), color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Proposed Solution
# ══════════════════════════════════════════════════════════════════════════════
def slide_03_solution(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "The Solution: One Platform. Full Intelligence Cycle.",
            "الحل: منصة واحدة. دورة ذكاء كاملة.",
            tag="الحل\nSOLUTION")
    _footer(s)

    # Intro statement
    _rect(s, Inches(0.35), Inches(1.42), SW - Inches(0.7), Inches(0.7), fill=NAVY)
    _txt(s,
         "GCC AI Intelligence Platform answers the three questions every GCC policymaker asks:  "
         "Where are we?  ·  Where are we heading?  ·  What should we do?",
         Inches(0.55), Inches(1.52), Inches(12.5), Inches(0.52),
         size=Pt(11.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    modules = [
        (NAVY,  "🌍", "GCC OVERVIEW",
         "نظرة إقليمية شاملة",
         "Six-country live dashboard · AI risk badges · YoY heatmap · GCC benchmarking"),
        (BLUE,  "🔍", "COUNTRY EXPLORER",
         "استكشاف معمق للدول",
         "KPIs · 15-year trend · strategic alerts · all-indicators snapshot"),
        (GREEN, "📈", "FORECAST CENTER",
         "مركز التنبؤ بالذكاء الاصطناعي",
         "6 AI models · auto-selection via cross-validation · trust score · confidence tiers"),
        (TEAL,  "🤖", "AI INSIGHTS",
         "رؤى الذكاء الاصطناعي",
         "Bilingual EN/AR reports · 8 risk profiles · alerts · policy recommendations"),
        (AMBER, "⚙️", "SCENARIO SIMULATOR",
         "محاكاة السيناريوهات",
         "8 policy presets · elasticity model · GCC-wide rank shift simulation"),
        (RED,   "🔬", "EXPLAINABILITY",
         "الشفافية والمصداقية",
         "Trust score 0–100 · driver intelligence · responsible AI documentation"),
    ]

    for i, (col, icon, title_en, title_ar, desc) in enumerate(modules):
        row, ci = divmod(i, 3)
        x = Inches(0.35 + ci * 4.32)
        y = Inches(2.3 + row * 2.45)
        _rect(s, x, y, Inches(4.1), Inches(0.65), fill=col)
        _txt(s, icon + "  " + title_en,
             x + Inches(0.15), y + Inches(0.08), Inches(3.8), Inches(0.3),
             size=Pt(10.5), bold=True, color=WHITE)
        _txt(s, title_ar,
             x + Inches(0.15), y + Inches(0.38), Inches(3.8), Inches(0.25),
             size=Pt(9), color=GOLD2)
        _rect(s, x, y + Inches(0.65), Inches(4.1), Inches(1.58),
              fill=WHITE, line=col, lw=Pt(2))
        _txt(s, desc,
             x + Inches(0.15), y + Inches(0.78), Inches(3.8), Inches(1.3),
             size=Pt(9.5), color=DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Implementation Mechanism
# ══════════════════════════════════════════════════════════════════════════════
def slide_04_implementation(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "How It Works: Data to Decision in One Flow",
            "آلية التنفيذ: من البيانات إلى القرار في خطوة واحدة",
            tag="آلية التنفيذ\nMECHANISM")
    _footer(s)

    # Pipeline boxes
    pipeline = [
        (GREEN, "🌐", "WORLD BANK\nOPEN DATA",
         "بيانات البنك\nالدولي",
         ["5 official indicators", "6 GCC nations", "2010–2024", "REST API + CSV offline cache"]),
        (NAVY,  "🗄", "DATA LAYER",
         "طبقة البيانات",
         ["Cleaning & validation", "Linear interpolation", "Offline-capable cache", "World Bank certified codes"]),
        (BLUE,  "🤖", "6-MODEL AI\nENGINE",
         "محرك الذكاء\nالاصطناعي",
         ["Naïve · Seasonal Naïve", "Moving Average · Drift", "SARIMAX (AIC auto)", "LightGBM (quantile CI)"]),
        (GOLD,  "✨", "INTELLIGENCE\nLAYER",
         "طبقة الذكاء\nالاستراتيجي",
         ["8 risk profiles", "4 alert types", "Bilingual EN + AR", "Trust score 0–100"]),
        (DARK,  "📊", "STREAMLIT\nDASHBOARD",
         "لوحة القيادة\nالتفاعلية",
         ["6 pages · Plotly charts", "Bilingual RTL UI", "4 export formats", "Demo flow guide"]),
    ]

    bw = Inches(2.35)
    aw = Inches(0.28)
    total = len(pipeline) * bw + (len(pipeline)-1) * aw
    sx = (SW - total) / 2

    for i, (col, icon, title_en, title_ar, items) in enumerate(pipeline):
        x = sx + i * (bw + aw)
        y = Inches(1.45)
        bh = Inches(4.85)
        # header
        _rect(s, x, y, bw, Inches(1.05), fill=col)
        _txt(s, icon + " " + title_en,
             x + Inches(0.1), y + Inches(0.04), bw - Inches(0.2), Inches(0.58),
             size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
        _txt(s, title_ar,
             x + Inches(0.1), y + Inches(0.64), bw - Inches(0.2), Inches(0.36),
             size=Pt(8.5), color=GOLD2, align=PP_ALIGN.CENTER)
        # body
        _rect(s, x, y + Inches(1.05), bw, bh - Inches(1.05),
              fill=WHITE, line=col, lw=Pt(1.8))
        _bullet_list(s, items, x + Inches(0.12), y + Inches(1.15),
                     bw - Inches(0.24), bh - Inches(1.2),
                     size=Pt(9.5), color=DARK, marker="▸ ")
        # arrow
        if i < len(pipeline) - 1:
            _txt(s, "→", x + bw + Inches(0.04), y + Inches(1.8), aw, Inches(0.42),
                 size=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Bottom highlight bar
    _rect(s, Inches(0.35), Inches(6.38), SW - Inches(0.7), Inches(0.65), fill=DARK)
    _txt(s,
         "✅  Offline-capable  ·  ✅  Automatic model selection via sMAPE CV  ·  "
         "✅  Bilingual EN+AR generated in parallel  ·  ✅  Production-ready one-click deploy",
         Inches(0.55), Inches(6.46), Inches(12.5), Inches(0.48),
         size=Pt(9.5), bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The AI Forecasting Engine (Results Part 1)
# ══════════════════════════════════════════════════════════════════════════════
def slide_05_forecasting(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "Six AI Models · Automatic Selection · Trust-Scored Results",
            "ستة نماذج ذكاء اصطناعي · اختيار تلقائي · نتائج مُقيَّمة بدرجة ثقة",
            tag="النتائج\nRESULTS")
    _footer(s)

    # Screenshot box — left
    _screenshot_box(s,
                    Inches(0.35), Inches(1.42), Inches(6.1), Inches(5.72),
                    "Forecast Center — Run Demo Forecast",
                    "أدخل صورة شاشة مركز التنبؤ هنا\n(بعد تشغيل التنبؤ التجريبي)",
                    "📸  Forecast Center page · Saudi Arabia · Youth Unemployment")

    # Right panel — model table + trust score
    _rect(s, Inches(6.75), Inches(1.42), Inches(6.25), Inches(5.72), fill=WHITE, line=NAVY, lw=Pt(1.5))
    _rect(s, Inches(6.75), Inches(1.42), Inches(6.25), Inches(0.5), fill=NAVY)
    _txt(s, "SIX-MODEL AI ENSEMBLE — AUTOMATIC SELECTION",
         Inches(6.92), Inches(1.48), Inches(6.0), Inches(0.37),
         size=Pt(10), bold=True, color=WHITE)

    models = [
        (MGRAY, "Naïve",           "Last-value persistence — baseline benchmark"),
        (MGRAY, "Seasonal Naïve",  "Seasonal pattern replication"),
        (MGRAY, "Moving Average",  "3-period local smoothing"),
        (MGRAY, "Drift",           "Linear trend extrapolation"),
        (BLUE,  "SARIMAX",         "Seasonal autoregression · AIC order selection"),
        (GOLD,  "LightGBM ★",      "Gradient boosting · quantile prediction intervals"),
    ]
    for i, (col, name, desc) in enumerate(models):
        y = Inches(2.08 + i * 0.73)
        _rect(s, Inches(6.85), y, Inches(0.18), Inches(0.52), fill=col)
        _txt(s, name, Inches(7.12), y + Inches(0.05), Inches(2.4), Inches(0.28),
             size=Pt(10.5), bold=True, color=DARK)
        _txt(s, desc, Inches(7.12), y + Inches(0.3), Inches(5.7), Inches(0.24),
             size=Pt(8.5), color=MGRAY)

    _rect(s, Inches(6.85), Inches(6.5), Inches(6.05), Inches(0.45), fill=GREEN)
    _txt(s,
         "★  Winner = lowest sMAPE on expanding-window holdout folds · sMAPE < 5% = A+ · < 12% = B",
         Inches(7.0), Inches(6.56), Inches(5.8), Inches(0.34),
         size=Pt(9), bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI Intelligence & Bilingual Reports (Results Part 2)
# ══════════════════════════════════════════════════════════════════════════════
def slide_06_intelligence(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "Ministry-Grade AI Reports — English & Arabic",
            "تقارير ذكاء اصطناعي بمستوى وزاري — إنجليزي وعربي",
            tag="النتائج\nRESULTS")
    _footer(s)

    # Large screenshot — left (AI Insights English)
    _screenshot_box(s,
                    Inches(0.35), Inches(1.42), Inches(6.1), Inches(2.65),
                    "AI Insights — Risk Panel + Strategic Alerts (English)",
                    "صورة شاشة: رؤى الذكاء الاصطناعي — لوحة المخاطر والتنبيهات",
                    "📸  AI Insights page · Saudi Arabia · English tab")

    # Screenshot — Arabic tab
    _screenshot_box(s,
                    Inches(0.35), Inches(4.25), Inches(6.1), Inches(2.88),
                    "AI Insights — Arabic Tab (العربية)",
                    "صورة شاشة: التقرير العربي للذكاء الاصطناعي",
                    "📸  AI Insights page · Arabic tab — RTL ministry-grade output")

    # Right panel — what the AI generates
    _rect(s, Inches(6.75), Inches(1.42), Inches(6.25), Inches(5.72), fill=WHITE, line=NAVY, lw=Pt(1.5))
    _rect(s, Inches(6.75), Inches(1.42), Inches(6.25), Inches(0.48), fill=TEAL)
    _txt(s, "🤖  WHAT THE AI GENERATES — ما يُنتجه الذكاء الاصطناعي",
         Inches(6.9), Inches(1.48), Inches(6.0), Inches(0.36),
         size=Pt(10), bold=True, color=WHITE)

    outputs = [
        (RED,   "⚡", "RISK CLASSIFICATION",   "تصنيف المخاطر",
         "8 dynamic risk labels auto-assigned\nبناءً على الاتجاه والتقلب والنطاق المستهدف"),
        (AMBER, "🔔", "STRATEGIC ALERTS",       "تنبيهات استراتيجية",
         "Up to 4 auto-generated alerts per forecast\nتنبيهات تلقائية: مستوى حرج، تحوّل سنوي، انحراف عن المتوسط"),
        (NAVY,  "📝", "EXECUTIVE BRIEF",        "الموجز التنفيذي",
         "Full bilingual narrative — EN + AR in parallel\nتقرير ثنائي اللغة جاهز للتنزيل خلال 60 ثانية"),
        (GREEN, "📊", "GCC COMPARISON",         "المقارنة الإقليمية",
         "Country vs GCC average — auto-generated\nتحليل مقارن: أين تقف الدولة بين أقرانها الخليجيين"),
    ]
    for i, (col, icon, en, ar, desc) in enumerate(outputs):
        y = Inches(2.08 + i * 1.2)
        _rect(s, Inches(6.85), y, Inches(0.12), Inches(0.95), fill=col)
        _txt(s, icon + "  " + en,
             Inches(7.07), y + Inches(0.04), Inches(5.7), Inches(0.35),
             size=Pt(10.5), bold=True, color=col)
        _txt(s, ar, Inches(7.07), y + Inches(0.38), Inches(5.7), Inches(0.25),
             size=Pt(9), color=MGRAY, italic=True)
        _txt(s, desc,
             Inches(7.07), y + Inches(0.62), Inches(5.7), Inches(0.5),
             size=Pt(8.5), color=DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Scenario Simulation (Results Part 3)
# ══════════════════════════════════════════════════════════════════════════════
def slide_07_scenario(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "Model the Future Before Committing to It",
            "محاكاة المستقبل قبل اتخاذ القرار",
            tag="النتائج\nRESULTS")
    _footer(s)

    # Formula bar
    _rect(s, Inches(0.35), Inches(1.42), SW - Inches(0.7), Inches(0.58), fill=DARK)
    _txt(s,
         "Scenario[t]  =  Baseline[t]  +  Σ (elasticity_i × delta_i)  ×  (t / horizon)   "
         "←  elasticity-based policy ramp",
         Inches(0.55), Inches(1.52), Inches(12.5), Inches(0.4),
         size=Pt(10.5), color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    # 8 scenario chips
    presets = [
        ("📊", "Baseline",            "الخط الأساسي",       MGRAY, WHITE),
        ("⚡", "Digital Acceleration","تسريع رقمي",          BLUE,  WHITE),
        ("📉", "Economic Slowdown",   "تباطؤ اقتصادي",       RED,   WHITE),
        ("🔥", "Inflation Stress",    "ضغط تضخمي",           AMBER, WHITE),
        ("🏗",  "Reform Expansion",   "توسع الإصلاح",        GREEN, WHITE),
        ("🌱", "Youth Recovery",      "تعافي الشباب",        TEAL,  WHITE),
        ("👥", "Population Surge",    "نمو سكاني",           NAVY,  WHITE),
        ("🚀", "High Growth GCC",     "نمو خليجي مرتفع",     DARK,  GOLD2),
    ]
    cw, ch = Inches(3.05), Inches(0.58)
    for i, (icon, en, ar, bg, tc) in enumerate(presets):
        row, ci = divmod(i, 4)
        x = Inches(0.35 + ci * (cw + Inches(0.1)))
        y = Inches(2.15 + row * (ch + Inches(0.14)))
        highlight = (en == "Reform Expansion")
        _rect(s, x, y, cw, ch, fill=GREEN if highlight else bg,
              line=GOLD if highlight else None, lw=Pt(2.5))
        _txt(s, icon + "  " + en,
             x + Inches(0.12), y + Inches(0.04), Inches(2.2), Inches(0.28),
             size=Pt(9.5), bold=True, color=GOLD if highlight else tc)
        _txt(s, ar,
             x + Inches(0.12), y + Inches(0.3), Inches(2.8), Inches(0.24),
             size=Pt(8.5), color=GOLD2 if highlight else RGBColor(0xCC,0xCC,0xCC))

    # Active scenario banner
    _rect(s, Inches(0.35), Inches(3.63), SW - Inches(0.7), Inches(0.42), fill=GREEN)
    _txt(s,
         "★  ACTIVE: Reform Expansion — Labour reform + Education investment + GDP growth → reduces youth unemployment",
         Inches(0.55), Inches(3.7), Inches(12.5), Inches(0.3),
         size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # KPI result cards
    kpis = [
        ("13.8%", "Baseline End",        "نهاية الخط الأساسي",    "Youth unemployment 2027 — no action",  MGRAY),
        ("11.3%", "Scenario End",        "نتيجة السيناريو",        "Under Reform Expansion preset",          GREEN),
        ("−2.5pp","Improvement",         "التحسن",                 "Percentage point delta from baseline",   TEAL),
        ("#4→#2", "GCC Rank Shift",      "تغير الترتيب الخليجي",   "Saudi Arabia rises among GCC peers",     GOLD),
    ]
    for i, (val, en, ar, desc, col) in enumerate(kpis):
        x = Inches(0.35 + i * 3.17)
        y = Inches(4.22)
        _rect(s, x, y, Inches(3.0), Inches(2.9), fill=WHITE, line=col, lw=Pt(2.5))
        _rect(s, x, y, Inches(3.0), Inches(0.38), fill=col)
        _txt(s, en,  x + Inches(0.1), y + Inches(0.06), Inches(2.8), Inches(0.28),
             size=Pt(8.5), bold=True, color=WHITE)
        _txt(s, val, x, y + Inches(0.5), Inches(3.0), Inches(0.95),
             size=Pt(36), bold=True, color=col, align=PP_ALIGN.CENTER)
        _txt(s, ar,  x + Inches(0.1), y + Inches(1.5), Inches(2.8), Inches(0.35),
             size=Pt(9.5), color=MGRAY, italic=True)
        _txt(s, desc,x + Inches(0.1), y + Inches(1.88), Inches(2.8), Inches(0.7),
             size=Pt(8.5), color=MGRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Explainability & Trust
# ══════════════════════════════════════════════════════════════════════════════
def slide_08_trust(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "AI That Earns Institutional Trust",
            "ذكاء اصطناعي يكسب الثقة المؤسسية",
            tag="النتائج\nRESULTS")
    _footer(s)

    # Trust score hero
    _rect(s, Inches(0.35), Inches(1.42), SW - Inches(0.7), Inches(1.18), fill=DARK)
    _rect(s, Inches(0.35), Inches(1.42), Inches(0.12), Inches(1.18), fill=GOLD)
    _txt(s, "TRUST SCORE",
         Inches(0.6), Inches(1.48), Inches(1.8), Inches(0.32),
         size=Pt(8.5), bold=True, color=GOLD)
    _txt(s, "درجة الثقة",
         Inches(0.6), Inches(1.78), Inches(1.8), Inches(0.28),
         size=Pt(8.5), color=GOLD2, italic=True)
    _rect(s, Inches(0.6), Inches(2.08), Inches(1.45), Inches(0.38), fill=AMBER)
    _txt(s, "74 / 100", Inches(0.6), Inches(2.08), Inches(1.45), Inches(0.38),
         size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s,
         "Elevated Uncertainty  ·  LightGBM  ·  sMAPE 10.4%  ·  Saudi Arabia · Youth Unemployment  ·  "
         "Tier C — Use alongside scenarios",
         Inches(2.2), Inches(1.6), Inches(10.9), Inches(0.5),
         size=Pt(10.5), color=WHITE)
    _txt(s,
         "World Bank Open Data  ·  Expanding-Window CV  ·  6-Model Ensemble  ·  Evidence-Based",
         Inches(2.2), Inches(2.08), Inches(10.9), Inches(0.38),
         size=Pt(9.5), color=GOLD2, italic=True)

    # Screenshot — Explainability page
    _screenshot_box(s,
                    Inches(0.35), Inches(2.75), Inches(5.8), Inches(4.35),
                    "Explainability — Trust Header + Driver Intelligence",
                    "صورة شاشة: لوحة الثقة وذكاء العوامل المحركة",
                    "📸  Explainability page · Saudi Arabia · after running forecast")

    # Right — 4 trust pillars
    pillars = [
        (GOLD,  "🎯", "TRUST SCORE 0–100",  "درجة ثقة 0–100",
         ["4 signals: sMAPE + horizon + interval width + volatility",
          "Single reliable metric for policymakers"]),
        (BLUE,  "🔑", "DRIVER INTELLIGENCE","ذكاء العوامل",
         ["LightGBM importances → policy-readable narratives",
          "Top drivers ranked with bilingual explanations"]),
        (NAVY,  "📐", "CONFIDENCE TIERS A–E","مستويات الثقة A–E",
         ["A (≥85): Strategic planning confidence",
          "C (≥55): Elevated uncertainty — use with scenarios"]),
        (TEAL,  "🛡",  "RESPONSIBLE AI",     "الذكاء الاصطناعي المسؤول",
         ["Data provenance · 6 limitation caveats",
          "Governance documentation · audit-ready"]),
    ]
    for i, (col, icon, en, ar, bullets) in enumerate(pillars):
        row, ci = divmod(i, 2)
        x = Inches(6.45 + ci * 3.42)
        y = Inches(2.75 + row * 2.25)
        _rect(s, x, y, Inches(3.22), Inches(2.08), fill=WHITE, line=col, lw=Pt(2))
        _rect(s, x, y, Inches(3.22), Inches(0.52), fill=col)
        _txt(s, icon + "  " + en,
             x + Inches(0.14), y + Inches(0.08), Inches(2.95), Inches(0.28),
             size=Pt(10), bold=True, color=WHITE)
        _txt(s, ar, x + Inches(0.14), y + Inches(0.34), Inches(2.95), Inches(0.2),
             size=Pt(8.5), color=GOLD2)
        _bullet_list(s, bullets, x + Inches(0.14), y + Inches(0.62),
                     Inches(2.95), Inches(1.35), size=Pt(9), color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Impact & GCC Alignment
# ══════════════════════════════════════════════════════════════════════════════
def slide_09_impact(prs):
    s = _blank(prs)
    _bg(s, LGRAY)
    _header(s,
            "Societal Impact & GCC Vision Alignment",
            "الأثر المجتمعي ومواءمة الرؤى الخليجية",
            tag="الأثر\nIMPACT")
    _footer(s)

    # Vision alignment table
    _rect(s, Inches(0.35), Inches(1.42), SW - Inches(0.7), Inches(0.45), fill=NAVY)
    _txt(s, "ALIGNED WITH ALL SIX GCC NATIONAL VISIONS  —  مواءمة مع جميع الرؤى الوطنية الخليجية الست",
         Inches(0.55), Inches(1.48), Inches(12.5), Inches(0.33),
         size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    visions = [
        ("🇸🇦", "Saudi Arabia", "Vision 2030",            "رؤية 2030",          GREEN),
        ("🇦🇪", "UAE",          "UAE Centennial 2071",    "مئوية الإمارات 2071", BLUE),
        ("🇶🇦", "Qatar",        "QNV 2030",               "رؤية قطر 2030",       RED),
        ("🇰🇼", "Kuwait",       "Kuwait Vision 2035",     "رؤية الكويت 2035",    AMBER),
        ("🇧🇭", "Bahrain",      "Economic Vision 2030",   "رؤية البحرين 2030",   TEAL),
        ("🇴🇲", "Oman",         "Oman Vision 2040",       "رؤية عُمان 2040",     NAVY),
    ]
    vw = (SW - Inches(0.7)) / 6
    for i, (flag, country, vision_en, vision_ar, col) in enumerate(visions):
        x = Inches(0.35) + i * vw
        _rect(s, x, Inches(1.9), vw, Inches(1.35), fill=col)
        _txt(s, flag, x, Inches(1.92), vw, Inches(0.52),
             size=Pt(20), align=PP_ALIGN.CENTER)
        _txt(s, country, x, Inches(2.4), vw, Inches(0.25),
             size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, vision_ar, x, Inches(2.65), vw, Inches(0.38),
             size=Pt(7.5), color=GOLD2, align=PP_ALIGN.CENTER, italic=True)

    # Three impact groups
    groups = [
        (NAVY, "🏛", "IMMEDIATE IMPACT",    "الأثر الفوري",
         ["Evidence-based forecasting replacing intuition",
          "Bilingual reports ready for cabinet briefings",
          "Strategic alerts before crises escalate",
          "Scenario testing before budget commitment"]),
        (GREEN,"🌍", "REGIONAL IMPACT",     "الأثر الإقليمي",
         ["Six-nation intelligence in one platform",
          "Policy learning: which countries lead, which lag",
          "Regional convergence / divergence tracking",
          "Comparative scenario analysis across all GCC"]),
        (TEAL, "🔮", "LONG-TERM IMPACT",   "الأثر طويل المدى",
         ["Institutional AI capacity within Gulf governments",
          "Audit-ready responsible AI documentation",
          "Scalable to MENA-wide coverage (15+ nations)",
          "REST API layer for government portal integration"]),
    ]
    for i, (col, icon, en, ar, bullets) in enumerate(groups):
        x = Inches(0.35 + i * 4.32)
        y = Inches(3.42)
        _rect(s, x, y, Inches(4.1), Inches(0.58), fill=col)
        _txt(s, icon + "  " + en,
             x + Inches(0.14), y + Inches(0.06), Inches(3.8), Inches(0.3),
             size=Pt(10.5), bold=True, color=WHITE)
        _txt(s, ar, x + Inches(0.14), y + Inches(0.36), Inches(3.8), Inches(0.22),
             size=Pt(8.5), color=GOLD2)
        _rect(s, x, y + Inches(0.58), Inches(4.1), Inches(3.55),
              fill=WHITE, line=col, lw=Pt(2))
        _bullet_list(s, bullets, x + Inches(0.18), y + Inches(0.72),
                     Inches(3.74), Inches(3.3), size=Pt(10), color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Creative Distinction & Conclusion
# ══════════════════════════════════════════════════════════════════════════════
def slide_10_distinction(prs):
    s = _blank(prs)
    _bg(s, NAVY2)
    _rect(s, Inches(0), Inches(0), Inches(0.22), SH, fill=GOLD)

    _txt(s, "نقاط التميز الإبداعي  ·  Creative Distinction",
         Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.65),
         size=Pt(28), bold=True, color=WHITE)
    _rect(s, Inches(0.4), Inches(0.9), Inches(8.0), Inches(0.07), fill=GOLD)
    _txt(s, "What makes this platform stand out from any other GCC labour market tool",
         Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.35),
         size=Pt(12), color=GOLD2, italic=True)

    distinctions = [
        (GOLD,  "01", "Ministry-Grade Bilingual AI",
         "الذكاء الاصطناعي الثنائي اللغة بمستوى وزاري",
         "Arabic executive reports generated in parallel — not translated. "
         "RTL-formatted, formal Modern Standard Arabic for GCC ministerial audiences."),
        (GREEN, "02", "AI That Explains Itself",
         "ذكاء اصطناعي يُفسّر نفسه",
         "Complete trust layer: 0-100 score, 5 confidence tiers, driver narratives. "
         "Every forecast comes with a full accountability trail."),
        (BLUE,  "03", "GCC-Wide Scenario Simulation",
         "محاكاة السيناريوهات على مستوى دول المجلس",
         "Test policy impact across ALL 6 nations simultaneously. "
         "See which country benefits most — before committing resources."),
        (TEAL,  "04", "Real Data. Not Synthetic.",
         "بيانات حقيقية، ليست مصطنعة",
         "Certified World Bank Open Data API — 5 official indicators, "
         "15 years, 6 nations. Offline-capable for field use."),
        (AMBER, "05", "Demo-Safe. Production-Ready.",
         "آمن للعرض. جاهز للإنتاج.",
         "One-click live URL — no setup, no install. "
         "Judges see a deployed production platform, not a prototype."),
    ]

    for i, (col, num, en, ar, desc) in enumerate(distinctions):
        row, ci = divmod(i, 3) if i < 3 else (1, i - 3)
        if i < 3:
            x = Inches(0.35 + i * 4.32)
            y = Inches(1.55)
            w, h = Inches(4.1), Inches(2.68)
        else:
            x = Inches(0.35 + (i - 3) * 6.55)
            y = Inches(4.42)
            w, h = Inches(6.25), Inches(2.68)

        _rect(s, x, y, w, h, fill=NAVY, line=col, lw=Pt(2.5))
        _rect(s, x, y, Inches(0.72), Inches(0.72), fill=col)
        _txt(s, num, x, y, Inches(0.72), Inches(0.72),
             size=Pt(22), bold=True, color=DARK, align=PP_ALIGN.CENTER)
        _txt(s, en,  x + Inches(0.82), y + Inches(0.06), w - Inches(0.95), Inches(0.35),
             size=Pt(10.5), bold=True, color=WHITE)
        _txt(s, ar,  x + Inches(0.82), y + Inches(0.4), w - Inches(0.95), Inches(0.28),
             size=Pt(9), color=GOLD2, italic=True)
        _rect(s, x + Inches(0.1), y + Inches(0.76), w - Inches(0.2), Inches(0.04), fill=col)
        _txt(s, desc, x + Inches(0.14), y + Inches(0.9), w - Inches(0.28), h - Inches(1.05),
             size=Pt(9.5), color=LGRAY, wrap=True)

    # Bottom gold quote bar
    _rect(s, Inches(0.22), Inches(7.12), SW - Inches(0.22), Inches(0.38), fill=GOLD)
    _txt(s,
         "🔗  Live:  competition-efvjqtcbnxcvt3te6fqkaw.streamlit.app   "
         "·   Data: World Bank Open Data API v2   "
         "·   Participant: Yahya Alwashahi",
         Inches(0.5), Inches(7.18), SW - Inches(0.6), Inches(0.28),
         size=Pt(9), bold=True, color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
def build():
    prs = _prs()
    print("Building slides...")
    slide_01_title(prs)          ; print("  ✓  1 — Cover / Title")
    slide_02_problem(prs)        ; print("  ✓  2 — Problem  (المشكلة)")
    slide_03_solution(prs)       ; print("  ✓  3 — Solution  (الحل المقترح)")
    slide_04_implementation(prs) ; print("  ✓  4 — Implementation  (آلية التنفيذ)")
    slide_05_forecasting(prs)    ; print("  ✓  5 — Results: AI Forecasting  (النتائج)")
    slide_06_intelligence(prs)   ; print("  ✓  6 — Results: AI Reports  (النتائج)")
    slide_07_scenario(prs)       ; print("  ✓  7 — Results: Scenario Sim  (النتائج)")
    slide_08_trust(prs)          ; print("  ✓  8 — Results: Trust & Explainability  (النتائج)")
    slide_09_impact(prs)         ; print("  ✓  9 — Impact  (الأثر)")
    slide_10_distinction(prs)    ; print("  ✓ 10 — Creative Distinction  (نقاط التميز)")

    os.makedirs("presentation", exist_ok=True)
    out = os.path.join("presentation", "GCC_AI_Intelligence_Platform.pptx")
    prs.save(out)
    print(f"\n✅  Saved: {out}")
    return out


if __name__ == "__main__":
    build()
